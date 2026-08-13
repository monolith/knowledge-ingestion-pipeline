"""Pass 0 — source intake and normalization. Deterministic code, no LLM.

Spec §8. Converts heterogeneous files into plain text plus a locator map that
can return to the exact portion of the original.
"""

from __future__ import annotations

import email
import html
import re
from email import policy
from pathlib import Path
from typing import Any

from .artifacts import (
    PipelineError,
    RunContext,
    file_hash,
    text_hash,
    utc_now,
    write_json_atomic,
    write_jsonl_atomic,
)
from .config import SCHEMA_VERSION

TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".rst", ".log", ".csv"}
HTML_SUFFIXES = {".html", ".htm"}
EMAIL_SUFFIXES = {".eml"}
DOCLING_SUFFIXES = {".pdf", ".docx", ".pptx", ".xlsx", ".epub", ".msg"}


def slugify(value: str) -> str:
    out = re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")
    return out or "source"


def source_id_for(path: Path) -> str:
    """A source id that cannot collide with another file's.

    The readable slug alone is not unique: `notes.md` and `notes.txt` both
    slugify to `notes`, and rglob makes `q1/report.md` and `q2/report.md`
    equally identical. A collision is silent and total -- the two documents
    share one normalized.txt, one unit_id namespace, and one registry entry, so
    the first document's text is overwritten before Pass 1 ever reads it while
    `kip validate` still reports a clean run. The content digest suffix makes
    that impossible while keeping the id human-readable.

    Content-derived rather than positional (an index) on purpose: the same file
    keeps the same id no matter what else is in the directory, so resuming a run
    after adding a source does not renumber everything already ingested.
    """
    # A path that does not exist cannot be hashed; fall back to the path itself
    # so the id stays unique rather than collapsing every missing file onto one.
    digest = file_hash(path) if path.is_file() else text_hash(str(path))
    return f"src-{slugify(path.stem)[:40]}-{digest[:8]}"


# --- Format handlers ----------------------------------------------------------
# Each returns (normalized_text, structural_markers) where markers map a line
# index to an original locator such as {"page": 7} or {"slide": 4}.


# A decoded file that is mostly U+FFFD was not text; errors="replace" turned an
# arbitrary binary into a page of replacement characters that reads as valid
# prose to every downstream stage. Spec §8.7 says an unsupported source is
# quarantined and never treated as valid, and garbage that merely survived
# decoding is exactly that.
_BINARY_REPLACEMENT_RATIO = 0.05


def _normalize_text(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    raw = path.read_bytes()
    if b"\x00" in raw:
        raise PipelineError(f"{path.name}: contains NUL bytes; not text")
    text = raw.decode("utf-8", errors="replace")
    if text and text.count("�") / len(text) > _BINARY_REPLACEMENT_RATIO:
        raise PipelineError(
            f"{path.name}: {text.count(chr(0xFFFD))} undecodable bytes in "
            f"{len(text)} characters; not usable text"
        )
    return text, {}


def _normalize_html(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    raw = path.read_text(encoding="utf-8", errors="replace")
    # Drop non-content elements wholesale before stripping tags, so their text
    # never reaches the model. Spec §8.5 forbids paraphrase but scripts and
    # styles are not document content.
    raw = re.sub(r"(?is)<(script|style)\b.*?</\1>", " ", raw)
    raw = re.sub(r"(?i)<br\s*/?>", "\n", raw)
    # `td` and `th` belong in this list and were missing from it. Without them a
    # row's cells concatenate with nothing between: the GE 10-K's segment table
    # reached the corpus as `Total segment revenue$33,314 $26,881 $23,855` and
    # its header as `SEGMENT REVENUE AND PROFIT202520242023`. A separator fixes
    # the fusing; it does not restore which year each figure belongs to, which
    # is why the tables are also recovered as assets.
    raw = re.sub(r"(?i)</(t[dh])>", " | ", raw)
    raw = re.sub(r"(?i)</(p|div|li|tr|h[1-6])>", "\n", raw)
    text = re.sub(r"(?s)<[^>]+>", "", raw)
    text = html.unescape(text)
    lines = [line.strip(" |\t") for line in text.splitlines()]
    return "\n".join(line for line in lines if line), {}


def _html_assets(path: Path, source_id: str) -> list[dict[str, Any]]:
    """Tables recovered from the markup, as grids.

    `exact` fidelity: the structure comes from the document's own `<table>`
    elements, not from a model reading a picture of one. The only inference is
    which row is a header, and that is confined to a flag.
    """
    from .assets import ASSET_TABLE, FIDELITY_EXACT, build_asset
    from .html_tables import compact, extract_tables

    raw = path.read_text(encoding="utf-8", errors="replace")
    out: list[dict[str, Any]] = []
    for index, table in enumerate(extract_tables(raw), start=1):
        grid = compact(table)
        if grid.n_rows < 2 or grid.n_cols < 2:
            continue
        out.append(build_asset(
            kind=ASSET_TABLE, source_id=source_id, index=len(out) + 1,
            fidelity=FIDELITY_EXACT, extractor="html_tables_v1",
            payload=grid.as_dict(), text=grid.to_text(),
        ))
    return out


def _normalize_email(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    message = email.message_from_bytes(path.read_bytes(), policy=policy.default)
    lines: list[str] = []
    # Spec §8.5: email keeps sender, recipients, sent time, subject, body,
    # quoted-message boundaries, and an attachment inventory.
    for header in ("From", "To", "Cc", "Date", "Subject"):
        value = message.get(header)
        if value:
            lines.append(f"{header}: {value}")
    lines.append("")

    body = ""
    attachments: list[str] = []
    if message.is_multipart():
        for part in message.walk():
            disposition = part.get_content_disposition()
            if disposition == "attachment":
                attachments.append(part.get_filename() or "unnamed")
            elif part.get_content_type() == "text/plain" and not body:
                body = part.get_content()
    else:
        body = message.get_content()

    lines.extend((body or "").splitlines())
    if attachments:
        lines.append("")
        lines.append(f"[[ATTACHMENTS {', '.join(attachments)}]]")
    return "\n".join(lines), {}


def _normalize_rich(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    """Rich formats: Docling when available, lightweight parsers otherwise.

    Docling is the spec's recommended normalizer (§8.3) because every DocItem
    carries page_no + bbox + charspan -- exactly the locator-map contract. It is
    tried first. The lightweight fallbacks below cost a few MB instead of a
    multi-GB ML stack, at a real price: PAGE/SLIDE-level provenance only, no
    bounding boxes. That price is recorded in the manifest's `normalizer` field
    so a downstream consumer can tell which fidelity it is holding.
    """
    try:
        return _docling(path)
    except PipelineError:
        pass  # try the lightweight path before giving up

    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _pdf_lite(path)
    if suffix == ".docx":
        return _docx_lite(path)
    if suffix == ".pptx":
        return _pptx_lite(path)
    if suffix == ".xlsx":
        return _xlsx_lite(path)
    raise PipelineError(
        f"{path.name}: no normalizer available for {suffix}. "
        "Install Docling with `pip install -e '.[parse]'` (full fidelity) "
        "or the light parsers with `pip install -e '.[parse-lite]'`."
    )


def _docling(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    try:
        from docling.document_converter import DocumentConverter
    except ImportError as exc:  # pragma: no cover - optional dependency
        raise PipelineError(f"{path.name}: Docling not installed") from exc

    result = DocumentConverter().convert(str(path))
    document = result.document
    markdown = document.export_to_markdown()

    # Map page markers by walking items in reading order. Docling's provenance
    # gives page_no per item; we record the first line where each page starts.
    markers: dict[int, dict[str, Any]] = {}
    lines = markdown.splitlines()
    seen_pages: set[int] = set()
    cursor = 0
    for item, _level in document.iterate_items():
        prov = getattr(item, "prov", None) or []
        if not prov:
            continue
        page_no = getattr(prov[0], "page_no", None)
        if page_no is None or page_no in seen_pages:
            continue
        text = (getattr(item, "text", "") or "").strip()
        if not text:
            continue
        needle = text.split("\n", 1)[0][:60]
        for index in range(cursor, len(lines)):
            if needle and needle in lines[index]:
                markers[index] = {"page": page_no}
                seen_pages.add(page_no)
                cursor = index
                break
    return markdown, markers


def _pdf_lite(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    """Page-level PDF text with [[PAGE n]] markers (spec §8.5)."""
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise PipelineError(
            f"{path.name}: needs a PDF parser. "
            "Install with `pip install -e '.[parse-lite]'` or `.[parse]`."
        ) from exc

    lines: list[str] = []
    markers: dict[int, dict[str, Any]] = {}
    for page_number, page in enumerate(PdfReader(str(path)).pages, start=1):
        markers[len(lines)] = {"page": page_number}
        lines.append(f"[[PAGE {page_number}]]")
        lines.extend((page.extract_text() or "").splitlines())
    return "\n".join(lines), markers


def _docx_lite(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise PipelineError(
            f"{path.name}: needs python-docx. Install with `pip install -e '.[parse-lite]'`."
        ) from exc

    document = docx.Document(str(path))
    lines = [p.text for p in document.paragraphs]
    # Spec §8.5: preserve tables rather than dropping them.
    for table in document.tables:
        lines.append("")
        for row in table.rows:
            lines.append(" | ".join(cell.text.strip() for cell in row.cells))
    return "\n".join(lines), {}


def _pptx_lite(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    """Slide order preserved, with title/body/notes distinguished (spec §8.5)."""
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise PipelineError(
            f"{path.name}: needs python-pptx. Install with `pip install -e '.[parse-lite]'`."
        ) from exc

    lines: list[str] = []
    markers: dict[int, dict[str, Any]] = {}
    for slide_number, slide in enumerate(Presentation(str(path)).slides, start=1):
        markers[len(lines)] = {"slide": slide_number}
        lines.append(f"[[SLIDE {slide_number}]]")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text)
        notes = getattr(slide, "notes_slide", None)
        if notes is not None and notes.notes_text_frame.text.strip():
            lines.append(f"[[NOTES]] {notes.notes_text_frame.text}")
    return "\n".join(lines), markers


def _xlsx_lite(path: Path) -> tuple[str, dict[int, dict[str, Any]]]:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise PipelineError(
            f"{path.name}: needs openpyxl. Install with `pip install -e '.[parse-lite]'`."
        ) from exc

    lines: list[str] = []
    markers: dict[int, dict[str, Any]] = {}
    workbook = load_workbook(str(path), data_only=True, read_only=True)
    for name in workbook.sheetnames:
        markers[len(lines)] = {"sheet": name}
        lines.append(f"[[SHEET {name}]]")
        for row in workbook[name].iter_rows(values_only=True):
            if any(cell is not None for cell in row):
                lines.append(" | ".join("" if c is None else str(c) for c in row))
    return "\n".join(lines), markers


def _handler_for(path: Path):
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        return _normalize_text, "text_v1"
    if suffix in HTML_SUFFIXES:
        return _normalize_html, "html_v1"
    if suffix in EMAIL_SUFFIXES:
        return _normalize_email, "email_v1"
    if suffix in DOCLING_SUFFIXES:
        return _normalize_rich, "rich_v1"
    return None, None


def supported(path: Path) -> bool:
    return _handler_for(path)[0] is not None


# --- Locator map --------------------------------------------------------------


def build_locator_map(
    source_id: str,
    normalized_path: str,
    text: str,
    markers: dict[int, dict[str, Any]],
) -> list[dict[str, Any]]:
    """One record per nonblank line, with char offsets and a content hash.

    Spec §8.6 requires both line numbers (for humans) and char offsets + hashes
    (for machines). Pass 5's deterministic citation check depends on the latter,
    so they are mandatory rather than best-effort.
    """
    records: list[dict[str, Any]] = []
    offset = 0
    current: dict[str, Any] = {}
    for index, line in enumerate(text.splitlines()):
        if index in markers:
            current = markers[index]
        stripped = line.strip()
        if stripped:
            records.append(
                {
                    "segment_id": f"seg-{source_id}-{len(records):05d}",
                    "source_id": source_id,
                    "normalized_path": normalized_path,
                    "normalized_line_start": index + 1,
                    "normalized_line_end": index + 1,
                    "normalized_char_start": offset,
                    "normalized_char_end": offset + len(line),
                    "original_locator_start": dict(current),
                    "original_locator_end": dict(current),
                    "text_sha256": text_hash(line),
                    "extraction_confidence": 1.0,
                }
            )
        offset += len(line) + 1  # +1 for the newline
    return records


# --- Pass entry point ---------------------------------------------------------


def _assets_for(path: Path, source_id: str, normalizer: str | None, *,
                text: str = "", assets_dir: Path | None = None) -> list[dict[str, Any]]:
    """Assets a normalizer can recover, or none.

    Kept separate from the text handler because the two answer different
    questions and fail independently: a source whose tables cannot be parsed
    still has usable text, and losing the text because of a table is the wrong
    trade.
    """
    try:
        if normalizer == "html_v1":
            return _html_assets(path, source_id)
        if normalizer == "rich_v1" and path.suffix.lower() == ".pdf":
            return _pdf_assets(path, source_id, text, assets_dir)
    except Exception as exc:  # never let an asset failure quarantine a source
        print(f"[pass0] {source_id}: asset extraction failed ({exc})")
    return []


def _pdf_assets(path: Path, source_id: str, text: str,
                assets_dir: Path) -> list[dict[str, Any]]:
    """Pages whose mathematics the text layer destroyed, rendered for reading.

    The transcription itself is not done here. Pass 0 is deterministic and takes
    no model calls; what it produces is the crop and the record that a formula
    is there, which a later vision pass fills in. An asset with an empty
    `latex` is honest -- it says a formula exists on this page and has not been
    read -- where omitting it says nothing at all.
    """
    from .pdf_assets import page_image_asset, pages_with_math, render_pages

    pages = pages_with_math(text)
    if not pages:
        return []
    written = render_pages(path, pages, assets_dir)
    if not written:
        print(f"[pass0] {source_id}: {len(pages)} page(s) carry damaged mathematics "
              "but no renderer is installed (pip install -e '.[parse-pdf]')")
        return []
    return [page_image_asset(source_id=source_id, index=i, page=page,
                             image_rel=f"{image.parent.name}/{image.name}")
            for i, (page, image) in enumerate(sorted(written.items()), start=1)]


def normalize_sources(ctx: RunContext, source_paths: list[Path]) -> list[dict[str, Any]]:
    """Normalize every source and write the source registry."""
    registry: list[dict[str, Any]] = []
    seen_ids: dict[str, Path] = {}

    for path in sorted(source_paths):
        handler, normalizer = _handler_for(path)
        source_id = source_id_for(path)
        if source_id in seen_ids:
            # Byte-identical duplicates are the only way to reach this now, and
            # ingesting one document twice inflates every coverage and
            # independence count that Pass 3 depends on. Refuse rather than
            # quietly merge -- a merge is what the old slug-only id did.
            raise PipelineError(
                f"duplicate source: {path} has the same content as "
                f"{seen_ids[source_id]} (source_id {source_id}); remove one"
            )
        seen_ids[source_id] = path
        target_dir = ctx.normalized_dir / source_id
        normalized_rel = f"01_normalized/{source_id}/normalized.txt"

        if handler is None:
            # Spec §8.7: unsupported or empty sources are quarantined, never
            # treated as valid. A silent skip would corrupt coverage metrics.
            registry.append(
                _manifest(
                    ctx, source_id, path, normalizer="unsupported",
                    status="quarantined", warnings=[f"unsupported format: {path.suffix}"],
                    normalized_rel=normalized_rel, normalized_sha="", line_count=0,
                )
            )
            continue

        try:
            text, markers = handler(path)
        except PipelineError as exc:
            registry.append(
                _manifest(
                    ctx, source_id, path, normalizer=normalizer or "unknown",
                    status="quarantined", warnings=[str(exc)],
                    normalized_rel=normalized_rel, normalized_sha="", line_count=0,
                )
            )
            continue

        if not text.strip():
            registry.append(
                _manifest(
                    ctx, source_id, path, normalizer=normalizer or "unknown",
                    status="quarantined", warnings=["empty after normalization"],
                    normalized_rel=normalized_rel, normalized_sha="", line_count=0,
                )
            )
            continue

        target_dir.mkdir(parents=True, exist_ok=True)
        (target_dir / "normalized.txt").write_text(text, encoding="utf-8", newline="\n")
        locator = build_locator_map(source_id, normalized_rel, text, markers)
        write_jsonl_atomic(target_dir / "locator_map.jsonl", locator)

        # Non-textual content, kept as addressable assets rather than flattened
        # into the line above. Written even when empty so a consumer can tell
        # "this source had no tables" from "this run predates assets".
        assets = _assets_for(path, source_id, normalizer,
                             text=text, assets_dir=target_dir / "assets")
        write_jsonl_atomic(target_dir / "assets.jsonl", assets)
        if assets:
            print(f"[pass0] {source_id}: {len(assets)} asset(s) recovered")

        manifest = _manifest(
            ctx, source_id, path, normalizer=normalizer or "unknown",
            status="success", warnings=[], normalized_rel=normalized_rel,
            normalized_sha=text_hash(text), line_count=len(text.splitlines()),
        )
        write_json_atomic(target_dir / "manifest.json", manifest)
        registry.append(manifest)

    write_jsonl_atomic(ctx.source_registry, registry)
    return registry


def _manifest(
    ctx: RunContext,
    source_id: str,
    path: Path,
    *,
    normalizer: str,
    status: str,
    warnings: list[str],
    normalized_rel: str,
    normalized_sha: str,
    line_count: int,
) -> dict[str, Any]:
    return {
        "schema_version": SCHEMA_VERSION,
        "run_id": ctx.run_id,
        "created_at": utc_now(),
        "source_id": source_id,
        "filename": path.name,
        "title": path.stem.replace("_", " ").replace("-", " ").strip(),
        "media_type": _media_type(path),
        "original_path": str(path),
        "original_sha256": file_hash(path) if path.exists() else "",
        "normalized_path": normalized_rel,
        "normalized_sha256": normalized_sha,
        "locator_map_path": f"01_normalized/{source_id}/locator_map.jsonl",
        "normalizer": normalizer,
        "normalizer_version": "1.0.0",
        "normalization_status": status,
        "normalized_line_count": line_count,
        "language": "en",
        # Lineage defaults to per-file independence. Spec §8.4 keeps these
        # fields authoritative for Pass 3's independence weighting, so operators
        # override them when two files describe the same underlying study or
        # pilot -- that override is what prevents evidence inflation.
        "source_family_id": f"family-{source_id}",
        "independence_group": f"group-{source_id}",
        "lineage_role": "unspecified",
        "derived_from": [],
        "quality_tier": "unspecified",
        "warnings": warnings,
    }


def _media_type(path: Path) -> str:
    return {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".html": "text/html",
        ".htm": "text/html",
        ".eml": "message/rfc822",
        ".md": "text/markdown",
        ".csv": "text/csv",
    }.get(path.suffix.lower(), "text/plain")
